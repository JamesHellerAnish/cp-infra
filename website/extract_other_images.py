import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

pptx_path = r'c:\Users\naben\OneDrive\Desktop\CP sytems\HYDROPOWER PRESENTATION.pptx'
out_dir = r'c:\Users\naben\OneDrive\Desktop\CP sytems\website\public\images\other_projects'
if not os.path.exists(out_dir):
    os.makedirs(out_dir)

prs = Presentation(pptx_path)

# Dictionary of slides to extract (0-indexed, so subtract 1 from visible slide numbers)
target_slides = {
    'tapovan': range(84, 89),
    'maneri_bhali': range(79, 84),
    'koteshwar': range(76, 79),
    'sewa': range(105, 107),
    'pipalkoti_grouting': range(69, 76),
    'general': [110]
}

metadata = []

for proj, slides in target_slides.items():
    img_idx = 1
    for slide_idx in slides:
        if slide_idx >= len(prs.slides): continue
        slide = prs.slides[slide_idx]
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                ext = image.ext
                img_bytes = image.blob
                filename = f"{proj}_{img_idx}.{ext}"
                with open(os.path.join(out_dir, filename), "wb") as f:
                    f.write(img_bytes)
                metadata.append(filename)
                img_idx += 1

print("Extracted files:", metadata)
print("Done extracting other project images.")
