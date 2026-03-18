import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

pptx_path = r'c:\Users\naben\OneDrive\Desktop\CP sytems\HYDROPOWER PRESENTATION.pptx'
out_dir = r'c:\Users\naben\OneDrive\Desktop\CP sytems\website\public\images\expansion'
if not os.path.exists(out_dir):
    os.makedirs(out_dir)

prs = Presentation(pptx_path)

target_slides = {
    'badrinath_rf_a': range(92, 99),
    'badrinath_civic': [99],
    'kedarnath_q1n': range(107, 110),
    'tehri_trt': [25],
    'tehri_drainage': range(26, 28),
    'tehri_surge': range(103, 106)
}

metadata = []

for proj, slides in target_slides.items():
    img_idx = 1
    for slide_idx in slides:
        if slide_idx >= len(prs.slides): continue
        slide = prs.slides[slide_idx]
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                ext = image.ext
                img_bytes = image.blob
                filename = f"{proj}_{img_idx}.{ext}"
                with open(os.path.join(out_dir, filename), "wb") as f:
                    f.write(img_bytes)
                metadata.append(filename)
                img_idx += 1

print("Extracted expansion files:", metadata)
print("Done extracting expansion images.")
