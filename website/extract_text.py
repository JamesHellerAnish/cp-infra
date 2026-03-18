import fitz
import os
from pptx import Presentation

pdf_path = r'c:\Users\naben\OneDrive\Desktop\CP sytems\C P SYSTEMS PVT LTD.pdf'
pptx_path = r'c:\Users\naben\OneDrive\Desktop\CP sytems\HYDROPOWER PRESENTATION.pptx'
out_path = r'c:\Users\naben\OneDrive\Desktop\CP sytems\website\extracted_text.txt'

with open(out_path, 'w', encoding='utf-8') as f:
    f.write("=== PDF CONTENT ===\n")
    try:
        doc = fitz.open(pdf_path)
        for i in range(len(doc)):
            f.write(f"--- Page {i+1} ---\n")
            f.write(doc[i].get_text("text") + "\n")
    except Exception as e:
        f.write(f"Error reading PDF: {e}\n")

    f.write("\n\n=== PPTX CONTENT ===\n")
    try:
        prs = Presentation(pptx_path)
        for i, slide in enumerate(prs.slides):
            f.write(f"--- Slide {i+1} ---\n")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    f.write(shape.text + "\n")
    except Exception as e:
        f.write(f"Error reading PPTX: {e}\n")

print("Extraction complete.")
